from typing import Dict, Any, List, Optional
import os
import json

from ..base_agent import BaseAgent
from ..models import Task, TaskResult, VerificationResult

# Lazy imports to avoid heavy startup if not used
# from docx import Document
# from openpyxl import Workbook
# from pptx import Presentation

class ProductivityAgent(BaseAgent):
    """
    Handles office productivity tasks: Word, Excel, PowerPoint
    """
    
    def __init__(self, agent_id: str = "productivity-001", config: Dict[str, Any] = None):
        super().__init__(agent_id, config or {})
        self.name = "Productivity Agent"
        self.category = "productivity"
        self.capabilities = [
            "create_word_document",
            "create_excel_spreadsheet",
            "create_powerpoint",
            "edit_document",
            "convert_formats"
        ]
        
    async def can_handle(self, task: Task) -> bool:
        """Check if task involves office documents"""
        keywords = ["document", "spreadsheet", "presentation", "pdf", 
                   "word", "excel", "powerpoint", "ppt", "docx", "xlsx", "doc", "xls"]
        return any(kw in task.description.lower() for kw in keywords)
    
    async def execute(self, task: Task) -> TaskResult:
        """Execute document creation/editing"""
        try:
            task_type = self._identify_task_type(task)
            print(f"ProductivityAgent executing task type: {task_type}")
            
            if task_type == "word_document":
                return await self._create_word_document(task)
            elif task_type == "excel_spreadsheet":
                return await self._create_spreadsheet(task)
            elif task_type == "powerpoint":
                return await self._create_presentation(task)
            else:
                return TaskResult(success=False, error=f"Unknown or unsupported document type: {task_type}")
                
        except Exception as e:
            import traceback
            traceback.print_exc()
            return TaskResult(success=False, error=str(e))
            
    def _identify_task_type(self, task: Task) -> str:
        desc = task.description.lower()
        if any(x in desc for x in ["word", "doc", "text file"]):
            return "word_document"
        elif any(x in desc for x in ["excel", "spreadsheet", "csv", "sheet"]):
            return "excel_spreadsheet"
        elif any(x in desc for x in ["powerpoint", "ppt", "slide", "presentation"]):
            return "powerpoint"
        return "unknown"
    
    # --- Word Handlers ---
    
    async def _create_word_document(self, task: Task) -> TaskResult:
        from docx import Document
        from docx.shared import Inches
        
        doc = Document()
        
        # Extract content
        title = task.params.get("title", "Document")
        content = task.params.get("content", "")
        
        doc.add_heading(title, 0)
        
        if isinstance(content, str):
            doc.add_paragraph(content)
        elif isinstance(content, list):
            for para in content:
                doc.add_paragraph(str(para))
                
        # Save
        filename = task.params.get("filename", f"document_{task.task_id[:8]}.docx")
        output_path = os.path.abspath(os.path.join("outputs", filename))
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        doc.save(output_path)
        
        return TaskResult(
            success=True,
            output_path=output_path,
            data={"filename": filename}
        )

    # --- Excel Handlers ---

    async def _create_spreadsheet(self, task: Task) -> TaskResult:
        from openpyxl import Workbook
        
        wb = Workbook()
        ws = wb.active
        ws.title = task.params.get("title", "Sheet1")
        
        data = task.params.get("data", [])
        
        # Expecting data as list of lists (rows)
        # or dict (headers + rows)
        if isinstance(data, list):
            for row in data:
                ws.append(row)
        
        filename = task.params.get("filename", f"sheet_{task.task_id[:8]}.xlsx")
        output_path = os.path.abspath(os.path.join("outputs", filename))
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        wb.save(output_path)
        
        return TaskResult(
            success=True,
            output_path=output_path,
            data={"filename": filename}
        )

    # --- PowerPoint Handlers ---

    async def _create_presentation(self, task: Task) -> TaskResult:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        
        # Title Slide
        title = task.params.get("title", "Presentation")
        subtitle = task.params.get("subtitle", "")
        
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        
        # Content Slides
        slides_data = task.params.get("slides", [])
        for slide_info in slides_data:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            if "title" in slide_info:
                slide.shapes.title.text = slide_info["title"]
            
            if "content" in slide_info:
                tf = slide.placeholders[1].text_frame
                tf.text = slide_info["content"]
        
        filename = task.params.get("filename", f"presentation_{task.task_id[:8]}.pptx")
        output_path = os.path.abspath(os.path.join("outputs", filename))
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        prs.save(output_path)
        
        return TaskResult(
            success=True,
            output_path=output_path,
            data={"filename": filename}
        )
